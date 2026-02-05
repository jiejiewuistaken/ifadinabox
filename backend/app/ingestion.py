from __future__ import annotations

import re
from dataclasses import dataclass
from pathlib import Path
from typing import Any
from uuid import uuid4

from pypdf import PdfReader
from docx import Document
from pptx import Presentation


@dataclass(frozen=True)
class Chunk:
    chunk_id: str
    doc_id: str
    source: str  # "user" | "internal"
    filename: str
    page: int | None
    text: str
    meta: dict[str, Any]


def _normalize_text(s: str) -> str:
    s = s.replace("\u00a0", " ")
    s = re.sub(r"[ \t]+", " ", s)
    s = re.sub(r"\n{3,}", "\n\n", s)
    return s.strip()


def extract_text_from_file(path: Path) -> list[tuple[str, dict[str, Any]]]:
    """
    Returns list of (text, meta) segments.
    For PDF we return per-page; for text/markdown we return a single segment.
    """
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        reader = PdfReader(str(path))
        out: list[tuple[str, dict[str, Any]]] = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            out.append((_normalize_text(text), {"page": i + 1}))
        return out
    if suffix == ".docx":
        doc = Document(str(path))
        parts: list[str] = []
        for p in doc.paragraphs:
            t = (p.text or "").strip()
            if t:
                parts.append(t)
        return [(_normalize_text("\n\n".join(parts)), {})]
    if suffix in (".txt", ".md"):
        return [(_normalize_text(path.read_text(encoding="utf-8", errors="ignore")), {})]
    if suffix == ".pptx":
        pres = Presentation(str(path))
        out: list[tuple[str, dict[str, Any]]] = []
        for i, slide in enumerate(pres.slides):
            parts: list[str] = []
            for shape in slide.shapes:
                text = ""
                if hasattr(shape, "text"):
                    text = shape.text or ""
                elif hasattr(shape, "text_frame") and shape.text_frame:
                    text = shape.text_frame.text or ""
                if text:
                    parts.append(text.strip())
            out.append((_normalize_text("\n\n".join(parts)), {"page": i + 1}))
        return out
    raise ValueError(f"Unsupported file type: {suffix}")


def chunk_text(
    text: str,
    *,
    chunk_size: int = 1200,
    overlap: int = 200,
) -> list[str]:
    """
    Simple, configurable chunker.
    - Prefer paragraph boundaries, then fallback to sliding window.
    """
    text = _normalize_text(text)
    if not text:
        return []

    paras = [p.strip() for p in text.split("\n\n") if p.strip()]
    chunks: list[str] = []
    cur = ""
    for p in paras:
        if len(cur) + len(p) + 2 <= chunk_size:
            cur = (cur + "\n\n" + p).strip()
        else:
            if cur:
                chunks.append(cur)
            cur = p
    if cur:
        chunks.append(cur)

    # Ensure max size via sliding window if needed
    final: list[str] = []
    for c in chunks:
        if len(c) <= chunk_size:
            final.append(c)
            continue
        start = 0
        while start < len(c):
            end = min(len(c), start + chunk_size)
            final.append(c[start:end].strip())
            if end == len(c):
                break
            start = max(0, end - overlap)
    return [x for x in final if x]


def build_chunks_for_file(
    path: Path, *, source: str, doc_id: str | None = None, meta: dict[str, Any] | None = None
) -> list[Chunk]:
    doc_id = doc_id or str(uuid4())
    segments = extract_text_from_file(path)
    filename = path.name
    chunk_meta = meta or {}
    out: list[Chunk] = []
    for seg_text, seg_meta in segments:
        page = seg_meta.get("page")
        for t in chunk_text(seg_text):
            out.append(
                Chunk(
                    chunk_id=str(uuid4()),
                    doc_id=doc_id,
                    source=source,
                    filename=filename,
                    page=page,
                    text=t,
                    meta=dict(chunk_meta),
                )
            )
    return out

